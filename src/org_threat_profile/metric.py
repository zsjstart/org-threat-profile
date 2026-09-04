import glob
import re
from enum import Enum, auto
from io import BytesIO
from typing import NamedTuple
from urllib.parse import urljoin

import ddgs
import ddgs.exceptions
import numpy as np
import pptx.shapes.autoshape
import pymupdf
import pymupdf4llm
import pytesseract
import requests
import tldextract
from bs4 import BeautifulSoup
from docx import Document as WordDocument
from markdownify import markdownify
from PIL import Image
from pptx import Presentation as PowerPointPresentation
from sentence_transformers import SentenceTransformer
from sentence_transformers.util import cos_sim
from thefuzz import fuzz
from tqdm import tqdm
from urlextract import URLExtract


class YamlKey:
    def __init__(self):
        self.keys = []
        self.single_indent = None

    def add_key(self, key_line: str):
        indent_match = re.match(r"^(\s+)\w", key_line)
        if indent_match:
            indent = indent_match.group(1)
            key = key_line.strip()
            if self.single_indent is None:
                self.single_indent = indent
                indent_level = 1
            else:
                indent_level = indent.count(self.single_indent)
        else:
            key = key_line
            indent_level = 0

        self.keys[indent_level:] = []
        key = key[: key.find(":")]  # RM semi-colon at the end
        self.keys.append(key)

    def __str__(self):
        return ":".join(self.keys)


class Fact(NamedTuple):
    key: str
    content: str


class SourceLocation(NamedTuple):
    title: str
    urls: list[str]


class ProcessedOutput(NamedTuple):
    stated_facts: dict[str, list[Fact]]
    inferred_facts: dict[str, list[Fact]]
    sources: dict[str, SourceLocation]


def process_output(model_output_file: str) -> ProcessedOutput:
    stated_facts = {}
    inferred_facts = {}
    used_sources = {}
    source_mode = False
    tracked_key = YamlKey()
    with open(model_output_file, "r") as f:
        for line in f:
            if not source_mode:
                if "# Source key:" in line:
                    source_mode = True
                    continue

                if re.match(r"^(\s*)\w+:", line):
                    tracked_key.add_key(line)

                if "(stated)" in line:
                    current_facts = stated_facts
                    fact_keyword = "(stated)"
                elif "(inferred)" in line:
                    current_facts = inferred_facts
                    fact_keyword = "(inferred)"
                else:
                    continue

                source_match = re.search(r"\[([^\]]+)\]", line)
                if source_match:
                    sources = source_match.group(1)
                    processed_line = line.replace(f"[{sources}]", "")
                    for source in sources.split(","):
                        source = source.strip()
                        processed_line = processed_line.replace(f" {fact_keyword}", "")
                        quoted_data = re.search(r'"([^"]*)"', processed_line)
                        if quoted_data:
                            processed_line = quoted_data.group(1)
                        processed_line = processed_line.strip()
                        fact = Fact(key=str(tracked_key), content=processed_line)
                        if current_facts.get(source) is None:
                            current_facts[source] = [fact]
                        else:
                            current_facts[source].append(fact)
            else:
                source_match = re.search(r"\[([^\]]+)\]", line)
                if source_match:
                    source = source_match.group(1)
                    title = line[line.find("]") + 1 : line.find("(")].strip()
                    urls = [
                        url.replace(")", "")
                        for url in URLExtract().find_urls(line)
                        if isinstance(url, str)
                    ]
                    used_sources[source] = SourceLocation(title=title, urls=urls)

    return ProcessedOutput(
        stated_facts=stated_facts, inferred_facts=inferred_facts, sources=used_sources
    )


def absolutise_urls(html: str, url: str) -> str:
    url_attributes = {
        "a": ["href"],
        "link": ["href"],
        "img": ["src", "srcset"],
        "script": ["src"],
        "iframe": ["src"],
        "source": ["src", "srcset"],
        "video": ["src", "poster"],
        "audio": ["src"],
        "form": ["action"],
    }
    soup = BeautifulSoup(html, "html.parser")
    base_url = "https://" + tldextract.extract(url).top_domain_under_public_suffix

    for tag_name, attributes in url_attributes.items():
        for tag in soup.find_all(tag_name):
            for attr in attributes:
                value = tag.get(attr)

                if not isinstance(value, str):
                    continue

                if attr == "srcset":
                    tag[attr] = ", ".join(
                        (
                            f"{urljoin(base_url, parts[0])}"
                            + (f" {parts[1]}" if len(parts) > 1 else "")
                        )
                        for parts in (part.strip().split() for part in value.split(","))
                    )
                else:
                    tag[attr] = urljoin(base_url, value)

    return str(soup)


def pdf_to_markdown(pdf_content: bytes) -> str:
    doc = pymupdf.open(stream=pdf_content, filetype="pdf")
    md = pymupdf4llm.to_markdown(doc, table_output="html")
    assert isinstance(md, str)
    return md


def html_to_markdown(html: str, url: str) -> str:
    html_with_abs_urls = absolutise_urls(html, url)
    md = markdownify(html_with_abs_urls)
    return md


def image_to_text(image_content: bytes) -> str:
    try:
        image = Image.open(BytesIO(image_content))
        return pytesseract.image_to_string(image)
    except pytesseract.TesseractNotFoundError:
        return "Could not perform OCR on the image, because the tesseract library is not installed."


def word_doc_to_text(doc_content: bytes) -> str:
    doc = WordDocument(BytesIO(doc_content))

    return "\n\n".join(
        paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()
    )


def pp_presentation_to_md(ppp_content: bytes) -> str:
    presentation = PowerPointPresentation(BytesIO(ppp_content))

    slides = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        text = "\n".join(
            shape.text
            for shape in slide.shapes
            if isinstance(shape, pptx.shapes.autoshape.Shape) and shape.text.strip()
        )
        if text.strip():
            slides.append(f"## Slide {slide_number}\n\n{text}")

    return "\n\n".join(slides)


class WebsiteContents(NamedTuple):
    content_type: str
    contents: str


class FailedRequest(Enum):
    Failed = auto()


def get_website_contents(url: str) -> WebsiteContents | FailedRequest:
    """
    Get the raw html of the website specified by the url.

    Args:
        url: URL of the website to get

    Returns:
        Markdown adapted from the raw contents of the website obtained via a GET request, no javascript or other DOM operations are executed by the website
    """
    try:
        if not url.startswith("http://") and not url.startswith("https://"):
            url = "https://" + url
        response = requests.get(url)
        response.raise_for_status()

        content_type = (
            response.headers.get("Content-Type", "").split(";", 1)[0].strip().lower()
        )
        match content_type:
            case "application/pdf":
                contents = pdf_to_markdown(response.content)
            case "text/html":
                contents = html_to_markdown(response.text, url)
            case s if s.startswith("image"):
                contents = image_to_text(response.content)
            case "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                contents = word_doc_to_text(response.content)
            case "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                contents = pp_presentation_to_md(response.content)
            case _:
                contents = response.text

        return WebsiteContents(content_type=content_type, contents=contents)

    except Exception:  # noqa: BLE001
        return FailedRequest.Failed


def perform_exact_search(
    title: str, limit: int = 10
) -> tuple[str, list[str]] | FailedRequest:
    try:
        matching_contents = ""
        matching_content_types = []
        searcher = ddgs.DDGS()
        results = searcher.text(
            query=title,
            max_results=limit,
            safesearch="off",
            region="wt-wt",
        )
        for result in results:
            if (
                title.lower() in result["title"].lower()
                or title.lower() in result["body"]
            ):
                website_contents = get_website_contents(result["href"])
                if website_contents is not FailedRequest.Failed:
                    matching_contents += website_contents.contents
                    matching_content_types.append(website_contents.content_type)
        return matching_contents, matching_content_types
    except ddgs.exceptions.DDGSException:
        return FailedRequest.Failed


def recursive_text_splitter(
    text: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
    separators: list[str] | None = None,
) -> list[str]:
    """
    Lightweight alternative to LangChain's RecursiveCharacterTextSplitter.

    Splitting priority:
        1. paragraph: "\n\n"
        2. line:      "\n"
        3. space:     " "
        4. character: ""

    Args:
        text: Input text.
        chunk_size: Maximum chunk length.
        chunk_overlap: Number of characters shared between chunks.
        separators: Custom separators, from coarse to fine.

    Returns:
        List of text chunks.
    """
    if chunk_overlap >= chunk_size:
        raise ValueError("chunk_overlap must be smaller than chunk_size")

    if separators is None:
        separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""]

    def split_recursive(text: str, separators: list[str]) -> list[str]:
        # Already small enough.
        if len(text) <= chunk_size:
            return [text]

        # Find the first separator that actually occurs.
        separator = separators[-1]
        remaining = []

        for i, sep in enumerate(separators):
            if sep == "" or sep in text:
                separator = sep
                remaining = separators[i + 1 :]
                break

        # Split into smaller pieces.
        if separator:
            pieces = text.split(separator)
        else:
            pieces = list(text)

        chunks = []
        current = []

        def flush():
            if current:
                chunks.append(separator.join(current))
                current.clear()

        for piece in pieces:
            if not piece:
                continue

            candidate = separator.join(current + [piece])

            if len(candidate) <= chunk_size:
                current.append(piece)
                continue

            # Current chunk is ready.
            flush()

            # This individual piece is still too large.
            if len(piece) > chunk_size and remaining:
                chunks.extend(split_recursive(piece, remaining))
            else:
                current.append(piece)

        flush()

        return chunks

    # First recursively split the document.
    pieces = split_recursive(text, separators)

    # Merge pieces while maintaining overlap.
    chunks = []
    current = ""

    for piece in pieces:
        if not current:
            current = piece
            continue

        candidate = current + "\n" + piece

        if len(candidate) <= chunk_size:
            current = candidate
            continue

        chunks.append(current)

        # Keep the tail of the previous chunk as overlap.
        overlap = current[-chunk_overlap:]

        # Avoid starting the next chunk with whitespace.
        overlap = overlap.lstrip()

        current = overlap + "\n" + piece

        # Extremely long piece: hard split as a safety measure.
        if len(current) > chunk_size:
            chunks.append(current[:chunk_size])
            current = current[chunk_size - chunk_overlap :]

    if current:
        chunks.append(current)

    return [chunk.strip() for chunk in chunks if chunk.strip()]


class SourceContent(NamedTuple):
    title: str
    urls: list[str]
    content: str
    content_types: list[str]
    chunks: np.typing.NDArray


class FailedSource(Enum):
    Failed = auto()


def crawl_sources(
    embedding_model: SentenceTransformer,
    sources: dict[str, SourceLocation],
) -> dict[str, SourceContent | FailedSource]:
    all_source_content = {}
    for source_name, source in sources.items():
        source_content = ""
        source_types = []
        if source.urls:
            for url in source.urls:
                website_contents = get_website_contents(url)
                if website_contents is not FailedRequest.Failed:
                    source_content += website_contents.contents
                    source_types.append(website_contents.content_type)
        else:
            exact_search_results = perform_exact_search(source.title)
            if exact_search_results is not FailedRequest.Failed:
                source_content, source_types = exact_search_results

        if len(source_content) == 0:
            all_source_content[source_name] = FailedSource.Failed
        else:
            all_source_content[source_name] = SourceContent(
                title=source.title,
                urls=source.urls,
                content=source_content.lower(),
                content_types=source_types,
                chunks=embedding_model.encode(recursive_text_splitter(source_content.lower())),
            )
    return all_source_content


class Score(NamedTuple):
    value: float
    total: int


def direct_score(
    embedding_model: SentenceTransformer,
    fact_content: str,
    source: SourceContent,
) -> float:
    if fact_content in source.content:
        return 1.0

    # Split source and get max sim
    if len(source.chunks) > 0:
        fact_emb = embedding_model.encode([fact_content])
        score = cos_sim(fact_emb, source.chunks).max().item()
        return score
    return 0.0


def stated_fact_score(
    embedding_model: SentenceTransformer,
    model_output: ProcessedOutput,
    sources: dict[str, SourceContent | FailedSource],
) -> Score:
    score_value = 0.0
    total_facts = 0

    for source_name, fact_list in model_output.stated_facts.items():
        total_facts += len(fact_list)
        source = sources[source_name]
        if source is FailedSource.Failed:
            continue

        for fact in fact_list:
            fact_content = fact.content.lower()
            score_value += direct_score(embedding_model, fact_content, source)

    return Score(value=score_value, total=total_facts)


def search_score(
    embedding_model: SentenceTransformer,
    topic: str,
    key: str,
    fact_content: str,
    t: int = 3,
    k: int = 10,
) -> float:
    try:
        searcher = ddgs.DDGS()
        results = searcher.text(
            query=f"{topic} {key} {fact_content}",
            max_results=k,
            safesearch="off",
            region="wt-wt",
        )
        search_scores = []
        for result in results:
            website_contents = get_website_contents(result["href"])
            if website_contents is not FailedRequest.Failed:
                search_contents = website_contents.contents
                source = SourceContent(
                    title=result["title"],
                    urls=[result["href"]],
                    content=search_contents,
                    content_types=[website_contents.content_type],
                    chunks=embedding_model.encode(recursive_text_splitter(search_contents)),
                )
                search_scores.append(
                    direct_score(embedding_model, fact_content, source)
                )
        search_scores.sort()
        if len(search_scores) > 0:
            return np.mean(search_scores[:-t]).item()
        return 0.0
    except ddgs.exceptions.DDGSException:
        return 0.0


def inferred_fact_score(
    embedding_model: SentenceTransformer,
    topic: str,
    model_output: ProcessedOutput,
    sources: dict[str, SourceContent | FailedSource],
) -> Score:
    score_value = 0.0
    total_facts = 0

    for source_name, fact_list in model_output.inferred_facts.items():
        total_facts += len(fact_list)
        source = sources[source_name]
        if source is FailedSource.Failed:
            continue

        for fact in fact_list:
            fact_content = fact.content.lower()
            direct_score_value = direct_score(embedding_model, fact_content, source)
            search_score_value = search_score(
                embedding_model, topic, fact.key.lower(), fact_content
            )
            score_value += max(direct_score_value, search_score_value)

    return Score(value=score_value, total=total_facts)


def hallucinated_score(sources: dict[str, SourceContent | FailedSource]) -> Score:
    """Find the number of sources that are hallucinated"""
    score_value = 0.0
    total_sources = len(sources)
    for source_content in sources.values():
        if source_content is FailedSource.Failed:
            score_value += 1
    return Score(value=score_value, total=total_sources)


def source_variety_score(sources: dict[str, SourceContent | FailedSource]) -> Score:
    """
    Measure the number of sources taken from unique locations,
    half points for sources that are different media type from the same location
    """
    total_sources = len(sources)
    used_sources: dict[str, set[str]] = {}  # source top level name -> content type
    score = 0.0
    for source in sources.values():
        if source is FailedSource.Failed:
            continue

        source_score = 0.0
        for url in source.urls:
            top_domain = tldextract.extract(url).top_domain_under_public_suffix
            if used_sources.get(top_domain) is None:
                source_score = 1.0
            elif any(ct not in used_sources[top_domain] for ct in source.content_types):
                source_score = max(source_score, 0.5)
            used_sources[top_domain] = used_sources.get(top_domain, set()) | set(
                source.content_types
            )
        score += source_score

    return Score(value=score, total=total_sources)


def score_output(
    embedding_model: SentenceTransformer, topic: str, model_output_file: str
) -> dict[str, Score]:
    model_output = process_output(model_output_file)
    sources = crawl_sources(embedding_model, model_output.sources)
    return {
        "stated_fact_score": stated_fact_score(embedding_model, model_output, sources),
        "inferred_fact_score": inferred_fact_score(
            embedding_model, topic, model_output, sources
        ),
        "hallucinated_sources": hallucinated_score(sources),
        "source_variety_score": source_variety_score(sources),
    }


def score_all_outputs(topic: str, output_folder: str) -> dict[str, dict[str, Score]]:
    embedding_model = SentenceTransformer("sentence-transformers/static-retrieval-mrl-en-v1")
    # embedding_model = SentenceTransformer("Qwen/Qwen3-Embedding-8B")
    scores = {}
    for output_file in (pbar := tqdm(glob.glob(f"{output_folder}/*.yaml"))):
        pbar.set_description_str(f"Processing {output_file}")
        scores[output_file[output_file.rfind("/") + 1 : output_file.rfind(".")]] = (
            score_output(embedding_model, topic, output_file)
        )
    return scores
